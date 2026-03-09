import os
import uuid
import tempfile
from flask import Flask, request, jsonify, send_file, render_template
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from werkzeug.utils import secure_filename

app = Flask(__name__)

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
ALLOWED_EXTENSIONS = {'.ppt', '.pptx'}

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static', 'logo.png')
LOGO_HEIGHT = Inches(0.45)
LOGO_MARGIN_TOP = Inches(0.2)
LOGO_MARGIN_RIGHT = Inches(0.3)

THEME = {
    'font_family': 'Poppins',
    'title_color': RGBColor(0x00, 0x27, 0xAB),
    'title_size': Pt(44),
    'heading_color': RGBColor(0x00, 0x27, 0xAB),
    'heading_size': Pt(32),
    'body_color': RGBColor(0x00, 0x00, 0x00),
    'body_size': Pt(24),
}

TITLE_PLACEHOLDER_INDICES = {0, 15}
SUBTITLE_PLACEHOLDER_INDICES = {1, 13}
BODY_PLACEHOLDER_INDICES = {2, 7, 14}

TITLE_SLIDE_LAYOUT_INDICES = {0, 5, 6}


def is_title_slide(slide):
    layout = slide.slide_layout
    layout_name = (layout.name or '').lower()
    if any(kw in layout_name for kw in ['title slide', 'title only', 'section']):
        return True
    try:
        idx = list(layout.slide_master.slide_layouts).index(layout)
        if idx in TITLE_SLIDE_LAYOUT_INDICES:
            return True
    except (ValueError, AttributeError):
        pass
    return False


def is_title_placeholder(placeholder):
    idx = placeholder.placeholder_format.idx
    ptype = placeholder.placeholder_format.type
    if ptype is not None and ptype in (1, 15):
        return True
    if idx in TITLE_PLACEHOLDER_INDICES:
        return True
    name = (placeholder.name or '').lower()
    if 'title' in name:
        return True
    return False


def is_subtitle_placeholder(placeholder):
    idx = placeholder.placeholder_format.idx
    ptype = placeholder.placeholder_format.type
    if ptype is not None and ptype == 2:
        return True
    if idx in SUBTITLE_PLACEHOLDER_INDICES:
        return True
    name = (placeholder.name or '').lower()
    if 'subtitle' in name:
        return True
    return False


def apply_theme_to_run(run, font_name, font_size, font_color, bold=False):
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold


def add_logo_to_slide(slide, prs):
    if not os.path.exists(LOGO_PATH):
        return

    slide_width = prs.slide_width
    logo_aspect = 474.0 / 256.0
    logo_h = LOGO_HEIGHT
    logo_w = int(logo_h * logo_aspect)

    left = slide_width - logo_w - LOGO_MARGIN_RIGHT
    top = LOGO_MARGIN_TOP

    slide.shapes.add_picture(LOGO_PATH, left, top, logo_w, logo_h)


def apply_theme(input_path, output_path):
    prs = Presentation(input_path)

    for slide_idx, slide in enumerate(prs.slides):
        title_slide = is_title_slide(slide)

        add_logo_to_slide(slide, prs)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            if shape.is_placeholder:
                ph = shape.placeholder_format

                if is_title_placeholder(shape):
                    if title_slide:
                        font_size = THEME['title_size']
                    else:
                        font_size = THEME['heading_size']
                    font_color = THEME['title_color']
                    bold = True
                elif is_subtitle_placeholder(shape):
                    font_size = THEME['body_size']
                    font_color = THEME['body_color']
                    bold = False
                else:
                    font_size = THEME['body_size']
                    font_color = THEME['body_color']
                    bold = False

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        apply_theme_to_run(
                            run, THEME['font_family'],
                            font_size, font_color, bold
                        )
            else:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        apply_theme_to_run(
                            run, THEME['font_family'],
                            THEME['body_size'], THEME['body_color'], False
                        )

    prs.save(output_path)


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/convert', methods=['POST'])
def convert():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    filename = secure_filename(file.filename)
    ext = os.path.splitext(filename)[1].lower()

    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({'error': 'Invalid file type. Only .ppt and .pptx files are allowed.'}), 400

    if ext == '.ppt':
        return jsonify({'error': 'Legacy .ppt format is not supported. Please convert to .pptx first.'}), 400

    temp_dir = tempfile.mkdtemp()
    input_path = os.path.join(temp_dir, f'input_{uuid.uuid4().hex}{ext}')
    output_path = os.path.join(temp_dir, f'themed_{filename}')

    try:
        file.save(input_path)

        file_size = os.path.getsize(input_path)
        if file_size > MAX_FILE_SIZE:
            return jsonify({'error': f'File too large. Maximum size is 50MB.'}), 400

        apply_theme(input_path, output_path)

        if os.path.exists(input_path):
            os.remove(input_path)

        response = send_file(
            output_path,
            as_attachment=True,
            download_name=f'themed_{filename}',
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        @response.call_on_close
        def cleanup():
            try:
                if os.path.exists(output_path):
                    os.remove(output_path)
                if os.path.exists(temp_dir):
                    os.rmdir(temp_dir)
            except OSError:
                pass

        return response

    except Exception as e:
        for f in [input_path, output_path]:
            try:
                if os.path.exists(f):
                    os.remove(f)
            except OSError:
                pass
        try:
            if os.path.exists(temp_dir):
                os.rmdir(temp_dir)
        except OSError:
            pass

        return jsonify({'error': f'Error processing file: {str(e)}'}), 500


if __name__ == '__main__':
    app.run(debug=True, port=5000)
